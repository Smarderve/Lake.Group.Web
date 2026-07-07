#!/usr/bin/env python3
"""Generate Lake Group stakeholder presentation (PPTX) — PDF-depth technical overview."""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "LAKE_GROUP_PRESENTATION.pptx"
GUIDE = ROOT / "docs" / "developer-guide.html"

# Brand palette
NAVY = RGBColor(0x0E, 0x1F, 0x5A)
BLUE = RGBColor(0x1D, 0x3E, 0xA8)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK_TEXT = RGBColor(0x1A, 0x23, 0x40)
MUTE = RGBColor(0x5A, 0x64, 0x78)
PAPER = RGBColor(0xF6, 0xF5, 0xF1)
LIGHT_BORDER = RGBColor(0xD8, 0xDE, 0xEC)
FOOTER_DARK = RGBColor(0xAA, 0xB4, 0xCC)
SUBTITLE_DARK = RGBColor(0xCC, 0xD4, 0xE8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# OOXML entrance animation preset IDs (PowerPoint 2016+)
EFFECT_PRESETS = {
    "appear": ("1", "0"),
    "fade": ("10", "0"),
    "fly": ("2", "4"),       # fly in from bottom
    "float": ("42", "0"),
    "wipe": ("22", "4"),
    "zoom": ("53", "16"),
}

TRANSITION_MAP = {
    "fade": "fade",
    "push": "push",
    "wipe": "wipe",
    "split": "split",
    "cover": "cover",
    "pull": "pull",
    "random": "randomBar",
}


# ---------------------------------------------------------------------------
# OOXML animation helpers
# ---------------------------------------------------------------------------

def _shape_spid(shape) -> str:
    return shape._element.xpath(".//p:cNvPr")[0].get("id")


def set_slide_transition(slide, transition_type: str = "fade", duration_ms: int = 500) -> None:
    """Inject p:transition on slide element (PowerPoint 2016+ compatible)."""
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    tag = TRANSITION_MAP.get(transition_type, transition_type)
    trans = etree.SubElement(sld, qn("p:transition"))
    trans.set("spd", "med")
    trans.set("advClick", "1")
    el = etree.SubElement(trans, qn(f"p:{tag}"))
    el.set("dur", str(duration_ms))
    if tag == "push":
        el.set("dir", "l")
    elif tag == "wipe":
        el.set("dir", "l")
    elif tag == "split":
        el.set("orient", "horz")
        el.set("dir", "out")


_anim_counter = 0


def _next_anim_id() -> int:
    global _anim_counter
    _anim_counter += 1
    return _anim_counter + 1


def add_entrance_animation(
    slide,
    shape,
    effect: str = "fade",
    delay_ms: int = 0,
    duration_ms: int = 400,
    click: bool = False,
) -> None:
    """Add shape entrance animation via p:timing / p:tnLst OOXML."""
    preset_id, preset_sub = EFFECT_PRESETS.get(effect, EFFECT_PRESETS["fade"])
    sid = _shape_spid(shape)
    sld = slide._element

    timing = sld.find(qn("p:timing"))
    if timing is None:
        timing = etree.SubElement(sld, qn("p:timing"))
    tn_lst = timing.find(qn("p:tnLst"))
    if tn_lst is None:
        tn_lst = etree.SubElement(timing, qn("p:tnLst"))

    root_par = tn_lst.find(qn("p:par"))
    if root_par is None:
        root_par = etree.SubElement(tn_lst, qn("p:par"))
        root_ctn = etree.SubElement(root_par, qn("p:cTn"))
        root_ctn.set("id", "1")
        root_ctn.set("dur", "indefinite")
        root_ctn.set("restart", "never")
        root_ctn.set("nodeType", "tmRoot")
        child_lst = etree.SubElement(root_ctn, qn("p:childTnLst"))
    else:
        root_ctn = root_par.find(qn("p:cTn"))
        child_lst = root_ctn.find(qn("p:childTnLst"))

    aid = _next_anim_id()
    par = etree.SubElement(child_lst, qn("p:par"))
    ctn = etree.SubElement(par, qn("p:cTn"))
    ctn.set("id", str(aid))
    ctn.set("fill", "hold")
    st_cond = etree.SubElement(ctn, qn("p:stCondLst"))
    cond = etree.SubElement(st_cond, qn("p:cond"))
    cond.set("delay", str(delay_ms))
    child2 = etree.SubElement(ctn, qn("p:childTnLst"))

    par2 = etree.SubElement(child2, qn("p:par"))
    ctn2 = etree.SubElement(par2, qn("p:cTn"))
    ctn2.set("id", str(aid + 10000))
    ctn2.set("fill", "hold")
    ctn2.set("presetID", preset_id)
    ctn2.set("presetClass", "entr")
    ctn2.set("presetSubtype", preset_sub)
    ctn2.set("nodeType", "clickEffect" if click else "withEffect")
    st2 = etree.SubElement(ctn2, qn("p:stCondLst"))
    cond2 = etree.SubElement(st2, qn("p:cond"))
    cond2.set("delay", "0")
    child3 = etree.SubElement(ctn2, qn("p:childTnLst"))

    set_el = etree.SubElement(child3, qn("p:set"))
    bhvr = etree.SubElement(set_el, qn("p:cBhvr"))
    bhvr.set("dur", str(duration_ms))
    bhvr.set("fill", "hold")
    ctn3 = etree.SubElement(bhvr, qn("p:cTn"))
    ctn3.set("id", str(aid + 20000))
    ctn3.set("dur", str(duration_ms))
    ctn3.set("fill", "hold")
    tgt = etree.SubElement(bhvr, qn("p:tgtEl"))
    sp_tgt = etree.SubElement(tgt, qn("p:spTgt"))
    sp_tgt.set("spid", sid)
    attr_lst = etree.SubElement(set_el, qn("p:attrNameLst"))
    attr = etree.SubElement(attr_lst, qn("p:attrName"))
    attr.text = "style.visibility"
    to_el = etree.SubElement(set_el, qn("p:to"))
    str_val = etree.SubElement(to_el, qn("p:strVal"))
    str_val.set("val", "visible")


def stagger_animations(slide, shapes, effect: str = "fade", base_delay: int = 80, duration: int = 350) -> None:
    for i, shape in enumerate(shapes):
        add_entrance_animation(slide, shape, effect=effect, delay_ms=i * base_delay, duration_ms=duration)


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------

class SlideBuilder:
    """Creates branded slides and applies transitions + entrance animations."""

    def __init__(self, prs: Presentation):
        self.prs = prs
        self._slide_idx = 0

    def _finish(self, slide, slide_kind: str, anim_shapes: list | None = None) -> None:
        transitions = {
            "cover": ("fade", 700),
            "section": ("push", 600),
            "content": ("fade", 500),
            "table": ("wipe", 550),
            "two_col": ("fade", 500),
            "stats": ("split", 550),
            "code": ("pull", 500),
            "closing": ("fade", 800),
        }
        t_type, t_dur = transitions.get(slide_kind, ("fade", 500))
        set_slide_transition(slide, t_type, t_dur)
        if anim_shapes:
            stagger_animations(slide, anim_shapes)
        self._slide_idx += 1

    def blank(self, dark: bool = False):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = NAVY if dark else WHITE
        return slide

    def add_corner_accent(self, slide, dark: bool = False):
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, SLIDE_W - Inches(0.18), 0, Inches(0.18), Inches(1.8))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD
        accent.line.fill.background()
        if not dark:
            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.42), Inches(0.55), Pt(3))
            rule.fill.solid()
            rule.fill.fore_color.rgb = GOLD
            rule.line.fill.background()

    def add_footer(self, slide, text: str = "lakeoilgroup.com  ·  July 2026", dark: bool = False):
        box = slide.shapes.add_textbox(Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.color.rgb = FOOTER_DARK if dark else MUTE
        p.font.name = "Calibri"

    def add_header_band(self, slide, title: str) -> object:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.28), Inches(11.5), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        return tb

    def cover(self, eyebrow: str, title: str, subtitle: str, stats: str):
        s = self.blank(dark=True)
        y = Inches(2.0)
        if eyebrow:
            eb = s.shapes.add_textbox(Inches(0.9), y, Inches(11), Inches(0.4))
            p = eb.text_frame.paragraphs[0]
            p.text = eyebrow.upper()
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.font.name = "Calibri"
            y += Inches(0.45)
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y, Inches(1.4), Pt(4))
        rule.fill.solid()
        rule.fill.fore_color.rgb = GOLD
        rule.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.9), y + Inches(0.2), Inches(11), Inches(2.0))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        sb = s.shapes.add_textbox(Inches(0.9), y + Inches(1.7), Inches(11), Inches(1.0))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = SUBTITLE_DARK
        sp.font.name = "Calibri"
        st = s.shapes.add_textbox(Inches(0.9), Inches(5.1), Inches(11), Inches(0.8))
        stp = st.text_frame.paragraphs[0]
        stp.text = stats
        stp.font.size = Pt(14)
        stp.font.color.rgb = FOOTER_DARK
        stp.font.name = "Calibri"
        self.add_footer(s, dark=True)
        self._finish(s, "cover", [tb, sb, st])

    def section(self, num: int, title: str, subtitle: str = ""):
        s = self.blank(dark=True)
        num_box = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(2), Inches(0.7))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{num:02d}"
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.font.name = "Calibri"
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(3.1), Inches(2.0), Pt(4))
        rule.fill.solid()
        rule.fill.fore_color.rgb = GOLD
        rule.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.9), Inches(3.4), Inches(11), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        shapes = [num_box, tb]
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11), Inches(0.8))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(16)
            sp.font.color.rgb = SUBTITLE_DARK
            sp.font.name = "Calibri"
            shapes.append(sb)
        self.add_footer(s, dark=True)
        self._finish(s, "section", shapes)

    def content(self, title: str, bullets: list[str], note: str | None = None, font_size: int = 18):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        body = s.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(4.9))
        tf = body.text_frame
        tf.word_wrap = True
        para_shapes = []
        for i, item in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = item
            para.level = 0
            para.font.size = Pt(font_size)
            para.font.color.rgb = INK_TEXT
            para.space_after = Pt(12)
            para.font.name = "Calibri"
        para_shapes.append(body)
        if note:
            nb = s.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11), Inches(0.5))
            np = nb.text_frame.paragraphs[0]
            np.text = note
            np.font.size = Pt(11)
            np.font.italic = True
            np.font.color.rgb = MUTE
            np.font.name = "Calibri"
            para_shapes.append(nb)
        self.add_footer(s)
        self._finish(s, "content", [header] + para_shapes)

    def two_col(self, title: str, left_title: str, left_items: list[str], right_title: str, right_items: list[str]):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        col_shapes = []

        def col(x, heading, items):
            hb = s.shapes.add_textbox(x, Inches(1.45), Inches(5.5), Inches(0.45))
            hp = hb.text_frame.paragraphs[0]
            hp.text = heading
            hp.font.size = Pt(15)
            hp.font.bold = True
            hp.font.color.rgb = BLUE
            hp.font.name = "Calibri"
            col_shapes.append(hb)
            box = s.shapes.add_textbox(x, Inches(1.95), Inches(5.5), Inches(4.6))
            tf = box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = f"• {item}"
                para.font.size = Pt(16)
                para.font.color.rgb = INK_TEXT
                para.space_after = Pt(8)
                para.font.name = "Calibri"
            col_shapes.append(box)

        col(Inches(0.9), left_title, left_items)
        col(Inches(6.8), right_title, right_items)
        self.add_footer(s)
        self._finish(s, "two_col", [header] + col_shapes)

    def stat_row(self, title: str, stats: list[tuple[str, str]]):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        n = len(stats)
        w = 11.0 / max(n, 1)
        cards = []
        for i, (num, label) in enumerate(stats):
            x = Inches(0.9 + i * w)
            card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(w - 0.25), Inches(3.6))
            card.fill.solid()
            card.fill.fore_color.rgb = PAPER
            card.line.color.rgb = LIGHT_BORDER
            cards.append(card)
            nb = s.shapes.add_textbox(x + Inches(0.12), Inches(2.35), Inches(w - 0.5), Inches(1.1))
            np = nb.text_frame.paragraphs[0]
            np.text = num
            np.font.size = Pt(34)
            np.font.bold = True
            np.font.color.rgb = BLUE
            np.alignment = PP_ALIGN.CENTER
            np.font.name = "Calibri"
            cards.append(nb)
            lb = s.shapes.add_textbox(x + Inches(0.12), Inches(3.5), Inches(w - 0.5), Inches(1.4))
            lp = lb.text_frame.paragraphs[0]
            lp.text = label
            lp.font.size = Pt(13)
            lp.font.color.rgb = MUTE
            lp.alignment = PP_ALIGN.CENTER
            lp.font.name = "Calibri"
            cards.append(lb)
        self.add_footer(s)
        self._finish(s, "stats", [header] + cards)

    def table(self, title: str, headers: list[str], rows: list[list[str]], col_widths: list[float] | None = None):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        n_cols = len(headers)
        n_rows = len(rows) + 1
        left, top = Inches(0.7), Inches(1.45)
        width, height = Inches(12.0), Inches(5.2)
        tbl = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table

        if col_widths:
            total = sum(col_widths)
            for ci, frac in enumerate(col_widths):
                tbl.columns[ci].width = int(width * frac / total)
        else:
            for ci in range(n_cols):
                tbl.columns[ci].width = int(width / n_cols)

        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"

        row_shapes = []
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFD)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = INK_TEXT
                    p.font.name = "Calibri"
        self.add_footer(s)
        self._finish(s, "table", [header, tbl._graphic_frame])

    def code_block(self, title: str, lines: str):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(12.0), Inches(5.2))
        box.fill.solid()
        box.fill.fore_color.rgb = PAPER
        box.line.color.rgb = LIGHT_BORDER
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.45), Pt(5), Inches(5.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.95), Inches(1.6), Inches(11.5), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines.strip().split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(11)
            para.font.name = "Consolas"
            para.font.color.rgb = INK_TEXT
            para.space_after = Pt(2)
        self.add_footer(s)
        self._finish(s, "code", [header, tb])

    def page_slide(self, filename: str, data: dict):
        s = self.blank()
        header = self.add_header_band(s, f"Page Reference — {filename}")
        self.add_corner_accent(s)
        y = Inches(1.42)
        blocks = []

        def add_label(text: str, bold: bool = False, color=None, size: int = 11):
            nonlocal y
            box = s.shapes.add_textbox(Inches(0.9), y, Inches(11.5), Inches(0.35))
            p = box.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color or INK_TEXT
            p.font.name = "Calibri"
            blocks.append(box)
            y += Inches(0.32)
            return box

        add_label(data["purpose"], size=13)
        add_label(f"~{data['lines']} lines  ·  CSS: {data['css']}", color=MUTE)
        add_label(f"JS: {data['js']}", color=MUTE)
        y += Inches(0.08)
        add_label("Sections", bold=True, color=BLUE)
        for sec in data["sections"][:4]:
            add_label(f"  • {sec}", size=10)
        add_label(f"Schema.org: {data['schema']}", color=MUTE, size=10)
        if data.get("unique"):
            add_label(f"Note: {data['unique']}", color=MUTE, size=10)
        self.add_footer(s)
        self._finish(s, "content", [header] + blocks)

    def agenda_multi(self, title: str, items: list[tuple[str, str]], cols: int = 2):
        s = self.blank()
        header = self.add_header_band(s, title)
        self.add_corner_accent(s)
        per_col = (len(items) + cols - 1) // cols
        col_w = 5.5
        boxes = []
        for ci in range(cols):
            chunk = items[ci * per_col : (ci + 1) * per_col]
            x = Inches(0.9 + ci * (col_w + 0.4))
            box = s.shapes.add_textbox(x, Inches(1.5), Inches(col_w), Inches(5.2))
            tf = box.text_frame
            tf.word_wrap = True
            for i, (num, label) in enumerate(chunk):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = f"{num}   {label}"
                para.font.size = Pt(13)
                para.font.color.rgb = INK_TEXT
                para.space_after = Pt(5)
                para.font.name = "Calibri"
            boxes.append(box)
        self.add_footer(s)
        self._finish(s, "content", [header] + boxes)

    def closing(self, eyebrow: str, title: str, subtitle: str):
        s = self.blank(dark=True)
        y = Inches(2.4)
        if eyebrow:
            eb = s.shapes.add_textbox(Inches(0.9), y, Inches(11), Inches(0.4))
            p = eb.text_frame.paragraphs[0]
            p.text = eyebrow.upper()
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.font.name = "Calibri"
            y += Inches(0.5)
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y, Inches(1.2), Pt(4))
        rule.fill.solid()
        rule.fill.fore_color.rgb = GOLD
        rule.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.9), y + Inches(0.2), Inches(11), Inches(1.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        sb = s.shapes.add_textbox(Inches(0.9), y + Inches(1.6), Inches(11), Inches(1.0))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = SUBTITLE_DARK
        sp.font.name = "Calibri"
        self.add_footer(s, "Lake Group — Powering East & Central Africa", dark=True)
        self._finish(s, "closing", [tb, sb])


# ---------------------------------------------------------------------------
# Content data (from DEVELOPER_GUIDE.pdf / developer-guide.html)
# ---------------------------------------------------------------------------

AGENDA_40 = [
    ("01", "Architecture Overview"),
    ("02", "Technology Stack"),
    ("03", "Production File Tree"),
    ("04", "Shared Page Anatomy"),
    ("05", "HTML Page Reference (29 pages)"),
    ("06", "JavaScript Modules"),
    ("07", "CSS Design Systems"),
    ("07b", "Meridian Component Catalog"),
    ("08", "3D Hero (hero-3d.js)"),
    ("08b", "3D Hero Deep Dive"),
    ("09", "Progressive Web App"),
    ("10", "Internationalization"),
    ("11", "Knowledge Assistant"),
    ("11b", "Assistant Deep Dive"),
    ("12", "News System"),
    ("13", "Africa Network Map"),
    ("14", "Build & Maintenance Scripts"),
    ("15", "SEO & Structured Data"),
    ("16", "Deployment Checklist"),
    ("17", "Accessibility"),
    ("18", "Appendices & Glossary"),
    ("19", "Troubleshooting"),
    ("20", "Firebase Hosting"),
    ("21", "Image Assets"),
    ("22", "i18n Key Conventions"),
    ("23", "Extended Page Docs"),
    ("24", "Extended JS Docs"),
    ("25", "Developer Workflows"),
    ("26", "Security Notes"),
    ("27", "Performance Budget"),
    ("28", "Version History"),
    ("29", "News Article Index"),
    ("30", "Subsidiaries & Page Mapping"),
    ("31", "Vendor Libraries"),
    ("32", "Glossary"),
    ("33", "Structured Data by Page"),
    ("34", "Cache Strategy Matrix"),
    ("35", "ADR Summary"),
    ("36", "Cross-Reference Index"),
    ("39", "Page Technical Essays"),
    ("40", "Build Scripts Encyclopedia"),
]

PAGES = [
    ("404.html", {"purpose": "Branded not-found page with navigation back to key sections.", "lines": 173,
     "css": "flagship.css, assistant.css", "js": "pwa.js, flagship-motion.js, assistant stack",
     "sections": [".error-page — 404 message and helpful links"],
     "schema": "None", "unique": "No i18n-content.js or site.js — minimal script set"}),
    ("about.html", {"purpose": "Company overview: founding story, values, subsidiary introduction. Meridian reference page.", "lines": 585,
     "css": "flagship.css, assistant.css", "js": "Standard stack + flagship-motion.js",
     "sections": [".page-hero", "#our-story-embed", ".fs-split sections", ".values-grid"],
     "schema": "BreadcrumbList", "unique": "First flagship-migrated page; template for Meridian rollouts"}),
    ("africa-network.html", {"purpose": "Interactive operations map of East & Central Africa with Leaflet satellite map.", "lines": 571,
     "css": "flagship.css, assistant.css, leaflet.css", "js": "leaflet.js, data_countries_africa.js, africa-network-map.js",
     "sections": [".page-hero", "#lake-africa-map", ".country-grid", ".map-legend"],
     "schema": "BreadcrumbList", "unique": "Leaflet loaded synchronously; requires network for tiles"}),
    ("careers.html", {"purpose": "Culture, benefits, CV submission form (mock).", "lines": 288,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Build Your Career Across Africa", "Drop Us Your CV"],
     "schema": "BreadcrumbList", "unique": "Migrated Meridian design"}),
    ("concrete.html", {"purpose": "GCCP ready-mix concrete and aggregate — Lugoba quarry, project portfolio.", "lines": 311,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Building East Africa's Future", "Our Plant in Action", "Specialist Concrete Portfolio"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "GCCP division page"}),
    ("contact.html", {"purpose": "HQ address, phone, email, regional offices, contact form.", "lines": 287,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["We're Happy to Hear From You"],
     "schema": "ContactPage, Organization, ContactPoint", "unique": "Primary conversion page"}),
    ("container-services.html", {"purpose": "AFICD inland container depots — Tanzania, Zambia, Mozambique.", "lines": 277,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["The Inland Gateway to East African Trade"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "AFICD & ACFS division"}),
    ("csr.html", {"purpose": "Corporate social responsibility programs and community investment.", "lines": 277,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Responsible Growth. Lasting Impact.", "How We Give Back"],
     "schema": "BreadcrumbList", "unique": "CSR programs and community stories"}),
    ("dashboard.html", {"purpose": "Mock client portal UI (sign-in, dashboard preview). Demo only.", "lines": 337,
     "css": "flagship.css, assistant.css", "js": "Standard stack",
     "sections": [".auth-panel", ".dashboard-preview"],
     "schema": "None", "unique": "robots.txt Disallow; not in sitemap or assistant KB"}),
    ("fleet.html", {"purpose": "700+ truck fleet specifications and capabilities.", "lines": 291,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Purpose-Built for Every Load"],
     "schema": "BreadcrumbList", "unique": "Lake Trans fleet showcase"}),
    ("fuel.html", {"purpose": "Lake Oil petroleum distribution — 85+ stations, bulk supply.", "lines": 284,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Keeping Africa Moving", "Our Competitive Edge"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "Flagship petroleum division page"}),
    ("gallery.html", {"purpose": "Masonry photo gallery with lightbox from operations and events.", "lines": 442,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Lake Group in Action"],
     "schema": "BreadcrumbList", "unique": "Largest interior page by line count"}),
    ("history.html", {"purpose": "Timeline from 2006 founding to present pan-African conglomerate.", "lines": 285,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Two Decades of Growth", "Be Part of the Next Chapter"],
     "schema": "BreadcrumbList", "unique": "Corporate timeline narrative"}),
    ("index.html", {"purpose": "Corporate homepage — 3D globe, divisions, founder, stats, CTAs.", "lines": 2931,
     "css": "theme.css, assistant.css", "js": "i18n stack, lazy hero-3d.bundle.js, motion.js",
     "sections": ["#fuel-experience / .experience-3d", ".hero", ".stats-grid", ".divisions-grid", ".founder-section"],
     "schema": "Organization, Person, PostalAddress", "unique": "Only page on theme.css; lazy-loads 3D bundle"}),
    ("investors.html", {"purpose": "Financial highlights, governance, currency-convertible figures.", "lines": 269,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["A Track Record of Vertical Growth", "What Drives Our Value"],
     "schema": "BreadcrumbList", "unique": "initCurrency() for FX converter"}),
    ("leadership.html", {"purpose": "Executive team, founder profile, subsidiary management.", "lines": 291,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Guided by Experience", "Ally Edha Awadh", "Group Management", "Companies Under Lake Group"],
     "schema": "BreadcrumbList", "unique": "Founder and executive profiles"}),
    ("logistics.html", {"purpose": "Lake Trans bulk liquid and dry cargo haulage.", "lines": 301,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["The Backbone of East Africa's Supply Chain", "Purpose-Built for Africa's Roads"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "Lake Trans division page"}),
    ("lpg.html", {"purpose": "Lake Gas LPG — cylinder sizes, bottling plants, clean cooking access.", "lines": 290,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Clean Energy for Every Kitchen", "Lake Gas Across Africa"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "Lake Gas division page"}),
    ("lubricants.html", {"purpose": "Lake Lubes — greases, industrial and automotive lubricants.", "lines": 267,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Engineered for Performance"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "Lake Lubes division page"}),
    ("media-center.html", {"purpose": "Press releases and media kit downloads.", "lines": 272,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Latest Announcements"],
     "schema": "BreadcrumbList", "unique": "Press-ready assets"}),
    ("news-article.html", {"purpose": "Single article view driven by ?id= query parameter.", "lines": 243,
     "css": "flagship.css, assistant.css", "js": "Standard + news-data.js + news.js",
     "sections": ["#article-root — title, date, gallery, YouTube embed"],
     "schema": "NewsArticle (dynamic)", "unique": "Dynamic JSON-LD; redirects if invalid id"}),
    ("news.html", {"purpose": "News listing with search, category filter, country filter.", "lines": 300,
     "css": "flagship.css, assistant.css", "js": "Standard + news-data.js + news.js",
     "sections": [".page-hero", ".news-filters", "#news-list"],
     "schema": "CollectionPage, BreadcrumbList", "unique": "20 articles from LAKE_NEWS array"}),
    ("offline.html", {"purpose": "PWA offline fallback when navigation fails without network.", "lines": 215,
     "css": "flagship.css, assistant.css", "js": "pwa.js, flagship-motion.js, assistant stack",
     "sections": [".offline-shell — links to precached pages"],
     "schema": "None", "unique": "Precached in sw.js; robots.txt Disallow"}),
    ("our-story.html", {"purpose": "Cinematic brand story — eight auto-advancing scenes with photography.", "lines": 619,
     "css": "flagship.css, assistant.css", "js": "i18n (sync), pwa.js, flagship-motion.js",
     "sections": [".ose-stage — scenes 1–8", ".story-stats overlay"],
     "schema": "None", "unique": "Not in main nav; keyboard/tap navigation"}),
    ("projects.html", {"purpose": "Major infrastructure and industrial project case studies.", "lines": 249,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Delivering East Africa's Development"],
     "schema": "BreadcrumbList", "unique": "Infrastructure case studies"}),
    ("services.html", {"purpose": "Hub page linking all eight business sectors.", "lines": 375,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["One Integrated Vision", "Ready to Work Together?"],
     "schema": "BreadcrumbList", "unique": "Division hub with 8 sector cards"}),
    ("station-locator.html", {"purpose": "Lake Oil retail station finder across Tanzania.", "lines": 310,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["85+ Stations Across Tanzania"],
     "schema": "BreadcrumbList", "unique": "Station locator for Lake Oil network"}),
    ("steel.html", {"purpose": "Lake Steel HS-CR rebars — 100,000 MT/year rolling mill.", "lines": 295,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["First in Tanzania. Built to Last.", "From Billets to Bars"],
     "schema": "Service, Organization, BreadcrumbList", "unique": "First HS-CR mill in Tanzania"}),
    ("sustainability.html", {"purpose": "Environmental stewardship, LPG clean cooking, ESG metrics.", "lines": 251,
     "css": "flagship.css, assistant.css", "js": "Standard deferred stack",
     "sections": ["Energy for Africa. Responsibly.", "Sustainability by the Numbers"],
     "schema": "BreadcrumbList", "unique": "ESG and environmental content"}),
]

TECH_STACK = [
    ("HTML5", "Page structure", "Zero build step; SEO-friendly; works everywhere"),
    ("CSS3 (inline + flagship/theme)", "Styling", "Critical CSS inlined for FCP; shared system for consistency"),
    ("Vanilla JavaScript (ES5/IIFE)", "Behaviour", "No framework runtime; small payload; file:// compatible"),
    ("Three.js (vendor)", "3D hero", "Industry-standard WebGL; bundled for non-module contexts"),
    ("Leaflet 1.x (vendor)", "Africa map", "Lightweight maps; Esri/OSM tile support"),
    ("FlexSearch (vendor)", "Assistant search", "Client-side full-text index; offline capable"),
    ("Firebase Hosting", "Deployment", "CDN, HTTPS; team uses Firebase CLI"),
    ("Service Worker", "PWA / offline", "Precache shell; stale-while-revalidate for assets"),
    ("Python 3", "Build tooling", "i18n extraction, font self-hosting, nav normalization"),
    ("Node.js", "KB build, SW lang, QA", "build_assistant_kb.js, add_* tag scripts"),
    ("esbuild", "3D bundle", "Inlines Three.js import into classic script"),
    ("Self-hosted fonts", "Typography", "Inter, Bebas Neue, Playfair — no Google Fonts CDN"),
]

COUNTRIES = [
    ("Tanzania", "TZ", "Dar es Salaam (HQ)", "Lake Oil, Lake Gas, Lake Steel, GCCP, Lake Trans, AFICD"),
    ("Kenya", "KE", "Nairobi", "Lake Oil Kenya, Mombasa port access"),
    ("Zambia", "ZM", "Lusaka / Ndola", "Lake Petroleum Zambia, AFICD"),
    ("Rwanda", "RW", "Kigali", "Lake Petroleum Rwanda, Lake Gas"),
    ("Burundi", "BI", "Bujumbura", "Burundi Petroleum"),
    ("DR Congo", "CD", "Lubumbashi", "DRC Petroleum, Lake Trans"),
    ("Ethiopia", "ET", "Addis Ababa", "Wadi Elsundus Petroleum"),
    ("Mozambique", "MZ", "Beira / Maputo", "Lake Oil LDA, AFICD"),
    ("UAE", "AE", "Dubai", "MERM Ready Mix, SAFF"),
]

CACHE_MATRIX = [
    ("navigate (HTML)", "Network-first → cache → offline.html", "All .html pages"),
    ("network-first-asset", "Network-first", "news-data.js, hero-3d.bundle.js"),
    ("cache-first-asset", "Cache-first", "/assets/fonts/, /assets/vendor/"),
    ("cache-first-image", "Cache-first + LRU (150)", "Most /assets/images/"),
    ("swr-image", "Stale-while-revalidate", "n-slider, lake-story-assets, products/"),
    ("swr-asset", "Stale-while-revalidate", "site.js, i18n, assistant, manifest"),
]

JS_MODULES = [
    ("site.js", "Shared behaviour: nav, reveal, counters, tabs, anchors, forms, currency"),
    ("i18n.js", "LakeI18n — EN/FR/SW; data-i18n attributes; localStorage persistence"),
    ("i18n-content.js", "Build artifact: ~1,442 translation keys as window global"),
    ("pwa.js", "SW registration, update toast, SKIP_WAITING, no surprise reloads"),
    ("motion.js", "theme.css partner — .reveal, nav hide/show, card tilt (index only)"),
    ("flagship-motion.js", "flagship.css partner — .fx-* choreography, magnetic CTAs"),
    ("hero-3d.js / .bundle.js", "Three.js globe — 39s loop, 9 sites, 4 facility callouts"),
    ("assistant.js", "FlexSearch retrieval, IndexedDB, no hallucination — verbatim passages"),
    ("assistant-kb.js", "Build artifact: knowledge base from i18n + curated facts"),
    ("news.js / news-data.js", "20 articles, filters, YouTube embeds, dynamic article pages"),
    ("africa-network-map.js", "Leaflet map — 22 markers, pipelines, country flyTo API"),
]

BUILD_SCRIPTS = [
    ("build_i18n_content.py", "Generate i18n-content.json + .js from master + translations"),
    ("build_master_en.py", "Extract English strings into _master_en.json"),
    ("build_assistant_kb.js", "Generate assistant-kb.js knowledge base"),
    ("build_hero_bundle.sh", "esbuild hero-3d.js → hero-3d.bundle.js"),
    ("normalize_nav.py", "Inject nav/footer/chat templates into all HTML pages"),
    ("add_seo_tags.js", "Inject OG, canonical, description meta tags"),
    ("add_pwa_tags.js", "Inject manifest link, theme-color, apple-touch-icon"),
    ("self_host_fonts.py", "Download and wire self-hosted font files"),
    ("_qa_check_links.js", "Broken link checker across all pages"),
    ("_qa_check_i18n.js", "i18n key coverage validation"),
    ("_globe_qa2.js", "3D hero automated frame capture QA"),
    ("generate_dev_pdf.sh", "PDF generation for developer guide"),
]


def build():
    global _anim_counter
    _anim_counter = 0

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    b = SlideBuilder(prs)

    # ── Cover ──────────────────────────────────────────────────────────────
    b.cover(
        "Lake Group",
        "Digital Platform\nTechnical Overview",
        "Architecture, systems & reference — aligned with DEVELOPER_GUIDE.pdf",
        "29 pages  ·  3 languages  ·  9 countries  ·  PWA-enabled  ·  www.lakeoilgroup.com",
    )

    # ── Agenda (40 PDF sections) ───────────────────────────────────────────
    b.agenda_multi("Agenda — Developer Guide Sections (1–20)", AGENDA_40[:20])
    b.agenda_multi("Agenda — Developer Guide Sections (21–40)", AGENDA_40[20:])

    b.stat_row(
        "Platform at a Glance",
        [
            ("29", "HTML pages"),
            ("3", "Languages\nEN · FR · SW"),
            ("1,442", "i18n keys"),
            ("9", "Verified\ncountries"),
            ("v10", "Service\nWorker"),
        ],
    )

    # ── §1 Architecture ────────────────────────────────────────────────────
    b.section(1, "Architecture Overview", "Static multi-page application — no SPA framework")
    b.code_block(
        "High-Level Architecture",
        """┌─────────────────────────────────────────────────────────────┐
│                     Browser (Client)                         │
│  index.html … sustainability.html  (29 pages)               │
│       ├── flagship.css OR theme.css  (design system)         │
│       ├── site.js + i18n.js + assistant.js                   │
│       └── page-specific: hero-3d, leaflet, news.js           │
├─────────────────────────────────────────────────────────────┤
│  Service Worker (sw.js) — precache, offline, cache routing │
│  manifest.webmanifest — installable PWA metadata             │
└─────────────────────────────────────────────────────────────┘
                              │
                              ▼
┌─────────────────────────────────────────────────────────────┐
│  Firebase Hosting  →  npm run deploy  →  global CDN + HTTPS   │
└─────────────────────────────────────────────────────────────┘""",
    )
    b.content(
        "Design Principles",
        [
            "Static-first: no server-side rendering; all copy ships in HTML + i18n-content.js",
            "file:// safe: critical data loads via <script> globals, not fetch()",
            "Progressive enhancement: pages render without JS; motion, assistant, 3D enhance when available",
            "Template-driven chrome: nav, footer, chat widget from scripts/templates/ via normalize_nav.py",
            "lake-3d/ Next.js project is NOT linked from production — live 3D is hero-3d.bundle.js only",
        ],
    )
    b.content(
        "Page Load Lifecycle",
        [
            "① Browser parses HTML; inline critical CSS paints first paint",
            "② External CSS (flagship/theme + assistant) loads",
            "③ Deferred scripts: i18n-content → i18n → site → flexsearch → assistant-kb → assistant → pwa → motion",
            "④ DOMContentLoaded: site.js inits nav, reveals, counters; LakeI18n applies stored language",
            "⑤ load: pwa.js registers service worker",
            "⑥ index.html only: IntersectionObserver lazy-loads hero-3d.bundle.js (+600px rootMargin)",
        ],
    )
    b.two_col(
        "Dual Design System Migration",
        "theme.css (Legacy)",
        [
            "index.html only",
            "Blue/yellow corporate, 4–10px radius",
            "motion.js (.reveal animations)",
            "Yellow divider bars",
        ],
        "flagship.css (Meridian)",
        [
            "28 interior pages migrated",
            "Ink/paper editorial, Bebas display",
            "flagship-motion.js (.fx-* choreography)",
            "Meridian tick + corner brackets",
        ],
    )

    # ── §2 Technology Stack ────────────────────────────────────────────────
    b.section(2, "Technology Stack", "Complete inventory with rationale")
    b.table("Technology Stack", ["Technology", "Role", "Why Chosen"], TECH_STACK, [1.2, 1.0, 2.3])

    # ── §3 File Tree ───────────────────────────────────────────────────────
    b.section(3, "Production File Tree", "Repository layout overview")
    b.code_block(
        "Key Production Files",
        """  ├── index.html, about.html … sustainability.html  (29 HTML pages)
  ├── assets/
  │   ├── flagship.css, theme.css, assistant.css
  │   ├── site.js, i18n.js, i18n-content.js, pwa.js
  │   ├── motion.js, flagship-motion.js, hero-3d.bundle.js
  │   ├── assistant.js, assistant-kb.js, news.js, news-data.js
  │   ├── africa-network-map.js, data_countries_africa.js
  │   ├── fonts/, images/, icons/pwa/, vendor/
  ├── scripts/  (build, QA, templates/)
  ├── sw.js, manifest.webmanifest, robots.txt, sitemap.xml
  └── docs/DEVELOPER_GUIDE.pdf, developer-guide.html""",
    )

    # ── §4 Shared Page Anatomy ─────────────────────────────────────────────
    b.section(4, "Shared Page Anatomy", "Standard HTML structure on every page")
    b.content(
        "Standard Page Structure",
        [
            "<head>: meta, canonical, OG tags, inline critical CSS, design system link",
            "<nav class='site-nav'>: desktop dropdowns, language switcher (EN/FR/SW)",
            "<div class='nav-mobile'>: accordion mobile menu",
            "<main>: page-hero + content sections with Meridian components",
            "<footer class='site-footer'>: four-column grid, social links, copyright",
            "<div id='chat-widget'>: assistant mount point (replaced by assistant.js)",
        ],
    )
    b.code_block(
        "Standard Script Load Order",
        """i18n-content.js → i18n.js → site.js → flexsearch.bundle.min.js
  → assistant-kb.js → assistant.js → pwa.js → flagship-motion.js

Homepage (index.html) replaces flagship pair with:
  theme.css + motion.js, adds lazy hero-3d.bundle.js via IntersectionObserver""",
    )

    # ── §5 HTML Pages (29) ─────────────────────────────────────────────────
    b.section(5, "HTML Page Reference", "Deep dive — all 29 root HTML pages")
    for filename, data in PAGES:
        b.page_slide(filename, data)

    # ── §6 JavaScript Modules ──────────────────────────────────────────────
    b.section(6, "JavaScript Modules", "assets/*.js — behaviour reference")
    for i in range(0, len(JS_MODULES), 2):
        chunk = JS_MODULES[i : i + 2]
        if len(chunk) == 2:
            b.two_col(
                "JavaScript Modules",
                chunk[0][0],
                [chunk[0][1]],
                chunk[1][0],
                [chunk[1][1]],
            )
        else:
            b.content("JavaScript Modules", [f"{chunk[0][0]}: {chunk[0][1]}"])
    b.table(
        "Vendor Libraries",
        ["File", "Version", "License", "Usage"],
        [
            ("vendor/three.module.min.js", "r16x", "MIT", "hero-3d.js WebGL rendering"),
            ("vendor/leaflet/leaflet.js", "1.x", "BSD-2", "africa-network.html map"),
            ("vendor/flexsearch/flexsearch.bundle.min.js", "0.7.x", "Apache-2.0", "assistant.js search"),
        ],
        [1.5, 0.6, 0.6, 1.8],
    )

    # ── §7 CSS Design Systems ──────────────────────────────────────────────
    b.section(7, "CSS Design Systems", "theme.css vs flagship.css (Meridian)")
    b.table(
        "Design System Comparison",
        ["Aspect", "theme.css (Legacy)", "flagship.css (Meridian)"],
        [
            ("Used on", "index.html only", "28 interior pages"),
            ("Design language", "Blue/yellow corporate", "Ink/paper editorial, Bebas display"),
            ("Dark sections", ".section-dark (navy)", ".fs-on-dark / --ink system"),
            ("Motion partner", "motion.js (.reveal)", "flagship-motion.js (.fx-*)"),
            ("Typography", "Inter + system", "Bebas Neue + Inter + Playfair quotes"),
            ("Tokens", "--blue, --yellow, --radius 4px", "--ink, --paper, --gold, --sp-* scale"),
        ],
        [1.0, 1.5, 1.5],
    )
    b.two_col(
        "Motion Systems Compared",
        "motion.js (index.html)",
        [
            "html.lg-motion gate class",
            ".reveal, .reveal-left, .reveal-scale",
            "Nav hide on scroll down",
            "3D card tilt (~3°) on pointer:fine",
        ],
        "flagship-motion.js (interior)",
        [
            "html.fs-motion gate class",
            ".fx-rise, .fx-clip, .fx-mask, .fx-wipe",
            "Scroll progress via [data-fx-scroll]",
            "Magnetic CTA hover; [data-fs-count] counters",
        ],
    )

    # ── §7b Meridian Components ────────────────────────────────────────────
    b.section(7, "Meridian Component Catalog", "flagship.css tokens & components")
    b.table(
        "Typography Classes",
        ["Class", "Font", "Usage"],
        [
            (".fs-hero-title", "Bebas Neue", "Page hero headlines, clamp 4–9.5rem"),
            (".fs-eyebrow", "Inter 600", "Tracked uppercase label with meridian tick"),
            (".fs-lede", "Inter", "Intro paragraph, larger size"),
            (".fs-serif-quote", "Playfair italic", "Pull quotes"),
            (".fs-stat-rail", "Bebas + Inter", "Horizontal KPI strip with baselines"),
        ],
        [1.2, 1.0, 2.3],
    )
    b.table(
        "Interactive Components",
        ["Class", "Behaviour"],
        [
            (".btn-primary", "Gold fill, notched corner, magnetic hover"),
            (".btn-blue", "Brand blue fill"),
            (".tab-nav / .tab-pane", "Tab switching via site.js initTabs()"),
            (".fx.fx-rise", "Translate Y + fade entrance"),
            (".fx.fx-clip", "Clip-path reveal for images/media"),
        ],
        [1.2, 2.3],
    )

    # ── §8 3D Hero ─────────────────────────────────────────────────────────
    b.section(8, "3D Global Operations Hero", "hero-3d.js — homepage cinematic experience")
    b.table(
        "Four Chapters (39-Second Loop)",
        ["Chapter", "Progress", "Duration", "Content"],
        [
            ("1 Planet", "0.00 – 0.30", "~10s", "Earth at distance; Africa rotates into view"),
            ("2 Footprint", "0.30 – 0.65", "~12s", "9 sites ignite from Dar HQ; route arcs draw"),
            ("3 Operations", "0.65 – 0.90", "~9s", "Facility callouts on leader lines"),
            ("4 Identity", "0.90 – 1.00", "~3s + 5s hold", "Logo finale overlay; fade; restart"),
        ],
        [0.8, 0.8, 0.7, 2.2],
    )
    b.table(
        "Verified Sites (SITES array)",
        ["Key", "City", "Country", "Notes"],
        [
            ("tz", "Dar es Salaam · HQ", "Tanzania", "Plot 49 Mikocheni — hub site"),
            ("ke", "Nairobi", "Kenya", "Kenya network base"),
            ("rw", "Kigali", "Rwanda", "Lake Petroleum Rwanda"),
            ("bu", "Bujumbura", "Burundi", "Burundi Petroleum"),
            ("cd", "Lubumbashi", "DR Congo", "DRC Petroleum"),
            ("zm", "Ndola", "Zambia", "Lake Petroleum Zambia"),
            ("et", "Addis Ababa", "Ethiopia", "Wadi Elsundus / WAS"),
            ("mz", "Beira", "Mozambique", "Lake Oil LDA + AFICD"),
            ("ae", "MERM · SAFF", "Dubai · UAE", "Middle East operations"),
        ],
        [0.5, 1.2, 0.9, 1.9],
    )
    b.content(
        "Facility Callouts & Textures",
        [
            "TANGA LPG TERMINAL — 3,000 MT storage, marine mooring buoy",
            "DAR ES SALAAM PORT — Vessel bunkering, 38M-litre depot",
            "LAKE STEEL · KIBAHA — HS-CR mill, 100,000 MT/yr",
            "GCCP · DAR ES SALAAM — Ready-mix, Lugoba quarry",
            "Textures: earth_color_2048.jpg, normal, specular, clouds from assets/images/planet/",
        ],
    )
    b.two_col(
        "3D Performance Optimizations",
        "Quality Tiers",
        [
            "isLow = width<768 OR cores<4 → fewer segments",
            "maybeDemoteQuality() after 90 frames >33ms",
            "IntersectionObserver pauses rAF off-screen",
            "Lazy bundle load — does not block LCP",
        ],
        "Memory & Fallbacks",
        [
            "InstancedMesh: 9 country markers, single draw call",
            "Canvas sprite labels generated once at init",
            "prefers-reduced-motion → static finale, no WebGL",
            "SW: network-first for hero-3d.bundle.js",
        ],
    )
    b.content(
        "3D Deep Dive — Implementation",
        [
            "boot() → checks reduced-motion → initHero3D() → is-ready class on success",
            "Globe radius GLOBE_R=3.4; lat/lon via colatitude + GLOBE_CENTER_LON=30°E",
            "Route arcs: TubeGeometry on quadratic Bezier; setDrawRange() animates drawing",
            "Label sprites: 3× supersampled canvas textures; SITE_NAMES[lang] on i18n event",
            "Camera: CAM_KEYS keyframes + smoothstep; mouse parallax ±6% (desktop only)",
            "Edit workflow: hero-3d.js → build_hero_bundle.sh → bump ?v=N → bump sw.js VERSION",
        ],
        note="QA: node scripts/_globe_qa2.js for automated frame capture",
    )

    # ── §9 PWA ─────────────────────────────────────────────────────────────
    b.section(9, "Progressive Web App", "Installable, offline-capable, graceful updates")
    b.table(
        "manifest.webmanifest Fields",
        ["Field", "Value"],
        [
            ("name / short_name", "Lake Group"),
            ("start_url", "./index.html"),
            ("display", "standalone"),
            ("theme_color", "#1D3EA8"),
            ("background_color", "#0E1F5A"),
            ("icons", "192, 512, maskable 192/512 in assets/icons/pwa/"),
        ],
        [1.2, 2.3],
    )
    b.table(
        "Four Cache Buckets (sw.js v10)",
        ["Cache Name", "Purpose"],
        [
            ("lake-precache-{VERSION}", "Install-time precache (shell, fonts, core JS/CSS)"),
            ("lake-pages-{VERSION}", "HTML navigations (network-first)"),
            ("lake-images-{VERSION}", "Images (cache-first, max 150 entries LRU)"),
            ("lake-assets-{VERSION}", "JS/CSS/fonts/vendor (stale-while-revalidate)"),
        ],
        [1.5, 2.0],
    )
    b.table("Fetch Routing Matrix", ["Route Class", "Strategy", "Examples"], CACHE_MATRIX, [1.2, 1.3, 1.8])
    b.content(
        "pwa.js Behaviour & Offline Pages",
        [
            "Registers SW on window load; resolves sw.js URL relative to pwa.js (subpath-safe)",
            "Update toast: navy background, gold border — 'Refresh' triggers SKIP_WAITING",
            "controllerchange reload only if user clicked Refresh (no surprise reloads)",
            "Silently no-ops on file:// or unsupported browsers",
            "offline.html: PWA fallback with links to precached about, services, contact",
            "404.html: branded not-found; assistant still works offline",
        ],
    )

    # ── §10 i18n ───────────────────────────────────────────────────────────
    b.section(10, "Internationalization", "English · French · Swahili — 1,442+ keys")
    b.code_block(
        "i18n Pipeline",
        """HTML pages (data-i18n attributes)
        │
        ▼
scripts/i18n_extract.py  →  discovers keys / orphans
        │
        ▼
scripts/build_master_en.py  →  scripts/_master_en.json
        │
        ▼
scripts/translation_dict.py  →  PHRASES_FR, PHRASES_SW
        │
        ▼
scripts/build_i18n_content.py  →  assets/i18n-content.js
        │
        ▼
assets/i18n.js  →  applies to DOM + language switch""",
    )
    b.content(
        "i18n Coverage & Attributes",
        [
            "Languages: English (en), French (fr), Swahili (sw) — Portuguese removed",
            "~1,442 translation keys across nav, footer, chat, homepage hero",
            "Attributes: data-i18n, data-i18n-html, data-i18n-placeholder, data-i18n-title, data-i18n-alt, data-i18n-aria",
            "Persists lake-lang in localStorage; dispatches lake-i18n-applied event",
            "hero-3d.js reads LakeI18n.current for SITE_NAMES country labels",
            "QA: node scripts/_qa_check_i18n.js validates key coverage",
        ],
    )

    # ── §11 Assistant ──────────────────────────────────────────────────────
    b.section(11, "Knowledge Assistant", "Retrieval-based — no hallucination")
    b.code_block(
        "Assistant Architecture",
        """scripts/build_assistant_kb.js
    ├── reads assets/i18n-content.json (page chunks)
    ├── merges CURATED_FACTS from _verified_lake_facts.md
    └── writes assets/assistant-kb.js

assets/assistant.js
    ├── FlexSearch index per language (en/fr/sw)
    ├── IndexedDB persistence (lake-assistant/messages)
    ├── Retrieves verbatim passages — never generates text
    └── "Read more →" links to source page""",
    )
    b.two_col(
        "Assistant Deep Dive",
        "Search Algorithm",
        [
            "Strip stopwords from query",
            "FlexSearch index with tokenized query",
            "Boost documents with f:1 (curated facts)",
            "Return top hit as answer with source URL",
        ],
        "Document Shape",
        [
            "{ id, t, s, u, k, f } — title, text, url, keywords, fact-boost",
            "Chunking: MIN_LEN=30, CHUNK_TARGET=340 chars",
            "Rebuild: node scripts/build_assistant_kb.js",
            "Works offline on 404.html and offline.html",
        ],
    )

    # ── §12 News ───────────────────────────────────────────────────────────
    b.section(12, "News System", "20 articles with search, filters, and video")
    b.content(
        "News Data Model & Rendering",
        [
            "window.LAKE_NEWS — array of articles: id, title, date, category, bannerImage, description[], images[], video",
            "news.html: renderNewsList(), filters by #news-search, #news-category, #news-country",
            "news-article.html: reads ?id= param; renders full article with YouTube embed",
            "Dynamic NewsArticle JSON-LD injected per article",
            "Images: assets/images/news/{id}/photo_*.jpg; SW network-first for news-data.js",
        ],
    )

    # ── §13 Africa Map ─────────────────────────────────────────────────────
    b.section(13, "Africa Network Map", "Leaflet — 22 markers, pipelines, country cards")
    b.two_col(
        "Map Dependencies & Tiles",
        "Load Order (synchronous)",
        [
            "leaflet.css + leaflet.js",
            "data_countries_africa.js → __LAKE_AFRICA_GEOJSON__",
            "africa-network-map.js",
        ],
        "Tile Layers",
        [
            "Hybrid (default): Esri World Imagery + Transportation",
            "Satellite, Terrain (OpenTopoMap), Streets (OSM)",
            "GeoJSON via script tag — file:// safe",
            "Tiles require network connection",
        ],
    )
    b.content(
        "Map Data Layers & API",
        [
            "countryLayer: GeoJSON fills for 8 ops countries + neighbours",
            "assetLayer: 22 circle markers by type (hq, fuel, port, container, industrial, logistics)",
            "pipelineLayer: TAZAMA pipeline, Northern corridor, Southern route polylines",
            "window.LakeAfricaMap = { flyToCountry(), filterAssets(), resetView() }",
            "Country cards call selectCountry(id, card) → map flyTo sync",
        ],
    )

    # ── §14 Build Scripts ──────────────────────────────────────────────────
    b.section(14, "Build & Maintenance Scripts", "Developer workflow automation")
    b.table("Build Scripts", ["Script", "Purpose"], BUILD_SCRIPTS, [1.5, 2.0])
    b.table(
        "QA & Automation Scripts",
        ["Script", "Purpose"],
        [
            ("_qa_matrix.js", "Cross-page QA matrix"),
            ("_qa_overflow_check.js", "Horizontal overflow detection"),
            ("_assistant_qa.js", "Assistant retrieval QA"),
            ("perf_pass.js", "Performance audit pass on pages"),
            ("inline_material_icons.py", "Inline Material Symbols SVGs"),
            ("add_theme_motion_tags.js", "Inject CSS/JS link tags"),
        ],
        [1.5, 2.0],
    )

    # ── §15 SEO ────────────────────────────────────────────────────────────
    b.section(15, "SEO & Structured Data", "Discoverable, shareable, Google-ready")
    b.content(
        "robots.txt & sitemap.xml",
        [
            "robots.txt: Allow / · Disallow /offline.html, /dashboard.html",
            "Sitemap: https://www.lakeoilgroup.com/sitemap.xml — 26 public URLs",
            "Priority 1.0 homepage, 0.8 news, 0.7 content pages; lastmod 2026-07-04",
            "Excluded: offline.html, dashboard.html, 404.html, our-story.html",
        ],
    )
    b.table(
        "Schema.org JSON-LD by Page Type",
        ["@type", "Pages"],
        [
            ("Organization", "index.html — company, founder, address"),
            ("Service", "fuel, lpg, steel, concrete, logistics, lubricants, container-services"),
            ("ContactPage", "contact.html — ContactPoint, Organization"),
            ("NewsArticle", "news-article.html (dynamic per article)"),
            ("BreadcrumbList", "Most interior pages"),
            ("CollectionPage", "news.html"),
        ],
        [1.0, 2.5],
    )
    b.content(
        "Open Graph & Meta Tags",
        [
            "Per-page: description, og:title, og:description, og:image (og-cover.jpg)",
            "twitter:card=summary_large_image on all major pages",
            "Canonical URLs: https://www.lakeoilgroup.com/{page}.html",
            "Maintained by scripts/add_seo_tags.js automated injector",
        ],
    )

    # ── §16 Deployment ─────────────────────────────────────────────────────
    b.section(16, "Deployment Checklist", "Firebase Hosting — npm run deploy")
    b.two_col(
        "Pre-Deploy Checklist",
        "Content Changes",
        [
            "build_master_en.py → build_i18n_content.py (if copy changed)",
            "node scripts/build_assistant_kb.js (if facts changed)",
            "bash scripts/build_hero_bundle.sh (if 3D changed)",
            "python3 scripts/normalize_nav.py (if templates changed)",
        ],
        "Version Bumps",
        [
            "Bump VERSION in sw.js for precache changes",
            "Bump ?v=N on hero-3d.bundle.js in index.html",
            "Update sitemap.xml lastmod for changed pages",
            "Run _qa_check_links.js and _qa_check_i18n.js",
        ],
    )
    b.content(
        "Post-Deploy Verification",
        [
            "Hard-refresh homepage — verify 3D hero loads after scroll",
            "Switch EN/FR/SW — verify nav and hero translate",
            "Open assistant — query 'fuel stations Tanzania'",
            "DevTools → Application → Service Workers — confirm new VERSION",
            "Test offline: load site, go offline, navigate to about.html",
            "Validate structured data: Google Rich Results Test",
        ],
        note="npm run serve → Firebase hosting emulator for local testing",
    )

    # ── §17 Accessibility ──────────────────────────────────────────────────
    b.section(17, "Accessibility", "WCAG-aligned foundations built in")
    b.two_col(
        "Reduced Motion & Semantics",
        "prefers-reduced-motion",
        [
            "motion.js/flagship-motion.js skip animation classes",
            "hero-3d.js shows static finale overlay, no WebGL",
            "CSS transitions disabled under reduced motion query",
        ],
        "Semantic HTML",
        [
            "<nav>, <main>, <footer>, <article> landmarks",
            "Heading hierarchy h1 → h2 → h3 per page",
            "Alt text; data-i18n-alt for translated alts",
        ],
    )
    b.content(
        "Keyboard, ARIA & Focus",
        [
            "Nav toggle: aria-label='Menu'; Language switcher: aria-label='Language'",
            "Assistant: aria-label on open/close; role='status' on PWA toast",
            "our-story.html: Space/→ keyboard navigation between scenes",
            "Flagship --gold-deep (#C79B00) for WCAG-safe gold on light surfaces",
            "Visible focus styles; assistant panel traps focus on mobile sheet",
        ],
    )

    # ── §26 Security ───────────────────────────────────────────────────────
    b.section(18, "Security & Privacy", "No tracking, local-only assistant data")
    b.content(
        "Security Notes",
        [
            "No user data sent to server — static site, mock forms only",
            "Assistant conversations stored locally in IndexedDB only",
            "SW never caches URLs with auth-like query params (token, apikey)",
            "No third-party analytics scripts in production HTML",
            "Map tiles contact Esri/OSM servers — disclose in privacy policy",
            "dashboard.html is demo UI only — disallow in robots.txt",
        ],
    )

    # ── §27 Performance ────────────────────────────────────────────────────
    b.content(
        "Performance Budget",
        [
            "index.html inline CSS ~80KB — first paint critical",
            "flagship.css ~45KB — render-blocking link",
            "hero-3d.bundle.js ~500KB — lazy on IntersectionObserver",
            "i18n-content.js ~200KB + assistant-kb.js ~150KB — defer",
            "Font files (woff2) ~300KB total — preload critical",
            "Target LCP: hero text before 3D loads; transform/opacity only for motion",
        ],
    )

    # ── §34 Countries ──────────────────────────────────────────────────────
    b.section(19, "Verified Operations Countries", "Source: scripts/_verified_lake_facts.md")
    b.table(
        "Nine Verified Countries",
        ["Country", "ISO", "Base City", "Entities"],
        COUNTRIES,
        [0.9, 0.5, 1.1, 2.0],
    )

    # ── §35 ADRs ───────────────────────────────────────────────────────────
    b.section(20, "Architectural Decision Records", "Key design choices documented")
    b.content(
        "ADR Summary",
        [
            "ADR-001 Static MPA: plain HTML + vanilla JS — zero framework runtime, independent cacheability",
            "ADR-002 Script-loaded globals: i18n, GeoJSON, KB, 3D bundle via <script> — file:// safe",
            "ADR-003 FlexSearch assistant: build-time KB + verbatim retrieval — no API costs, no hallucination",
            "ADR-004 Dual design system: incremental Meridian rollout — both precached until index.html migrates",
            "ADR-005 Lazy-load 3D: IntersectionObserver +600px — ~500KB bundle must not block LCP",
        ],
        note="Full ADR text in DEVELOPER_GUIDE.pdf §36",
    )

    # ── §20 Firebase ───────────────────────────────────────────────────────
    b.content(
        "Firebase Hosting",
        [
            "Production URL: https://www.lakeoilgroup.com/",
            "Deploy: npm run deploy → firebase deploy --only hosting",
            "Local: npm run serve → Firebase hosting emulator",
            "Global CDN, HTTPS by default, clean cache headers",
            "firebase.json + .firebaserc configured locally (not committed)",
        ],
    )

    b.table(
        "Migration & Version Status",
        ["Component", "Status", "Notes"],
        [
            ("Flagship CSS rollout", "28/29 pages", "index.html remains on theme.css"),
            ("i18n architecture", "Complete", "~1,442 keys; partial long-tail coverage"),
            ("PWA / Service Worker", "v10", "Bump VERSION per deploy"),
            ("Knowledge Assistant", "Live", "Replaces legacy chat"),
            ("hero-3d bundle", "v9 query param", "file:// safe; lazy-loaded"),
            ("Portuguese (PT)", "Removed", "Replaced with Swahili (SW)"),
        ],
        [1.2, 0.9, 2.4],
    )

    b.table(
        "Subsidiaries & Page Mapping",
        ["Subsidiary", "Sector", "Page"],
        [
            ("Lake Oil", "Fuel & Petroleum", "fuel.html"),
            ("Lake Gas", "LPG", "lpg.html"),
            ("Lake Lubes", "Lubricants", "lubricants.html"),
            ("Lake Steel", "Steel & Rebars", "steel.html"),
            ("GCCP", "Concrete & Aggregate", "concrete.html"),
            ("Lake Trans", "Logistics & Haulage", "logistics.html, fleet.html"),
            ("AFICD / ACFS", "Container Services", "container-services.html"),
            ("MERM Dubai", "Ready Mix (UAE)", "Referenced in 3D hero (ae)"),
        ],
        [1.0, 1.2, 1.3],
    )

    b.content(
        "Brand Colour Reference",
        [
            "Navy #0E1F5A — headers, ink bands, PWA toast",
            "Blue #1D3EA8 — primary brand, theme-color, buttons",
            "Gold #FFD700 — accents, CTAs, meridian ticks (WCAG: --gold-deep #C79B00 on light)",
            "Paper #F6F5F1 — flagship light sections",
            "Ink #070C1E — flagship dark bands",
        ],
    )

    # ── Closing ────────────────────────────────────────────────────────────
    b.content(
        "Key Takeaways for Stakeholders",
        [
            "Premium verified corporate platform — 29 pages, 3 languages, 9 countries",
            "3D homepage tells the Africa operations story with verified data only",
            "Installable PWA — works offline, updates gracefully, no app store",
            "SEO-ready — sitemap, Schema.org, Open Graph on every public page",
            "Knowledge assistant — instant auditable answers, works offline",
            "Meridian design system — precision-engineering editorial across 28 pages",
        ],
    )
    b.content(
        "Documentation & Maintenance",
        [
            "DEVELOPER_GUIDE.pdf — full technical reference (80+ pages, 40 sections)",
            "docs/developer-guide.html — browsable HTML version with deep links",
            "Content verified against scripts/_verified_lake_facts.md",
            "Build scripts for translations, assistant KB, 3D bundle, nav templates",
            "Designed for long-term maintainability and incremental enhancement",
        ],
        note="Repository: github.com/Smarderve/Lake.Group.Web",
    )
    b.closing(
        "Thank You",
        "Questions &\nDiscussion",
        "Full reference: DEVELOPER_GUIDE.pdf  ·  lakeoilgroup.com",
    )

    prs.save(OUT)
    size = OUT.stat().st_size
    count = len(prs.slides)
    print(f"Wrote {OUT}")
    print(f"  Slides: {count}")
    print(f"  Size:   {size:,} bytes ({size / 1024:.1f} KB)")
    print(f"  Animations: OOXML transitions on all slides + entrance animations on headers/bullets/cards/tables")


if __name__ == "__main__":
    build()
